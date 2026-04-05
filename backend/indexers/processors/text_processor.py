# -*- coding: utf-8 -*-
"""Text, DOCX, and PPTX processing for ChunkyLink document indexing."""
import logging
from pathlib import Path
from typing import List, Dict, Any

from backend.indexers.utils.text_processing import (
    normalize_text, summarize_for_router, sha8, split_with_overlap, should_deduplicate
)
from backend.indexers.utils.ocr_processor import process_docx_images, process_pptx_images


def build_for_txt(txt_path: Path, cfg: Dict[str, Any]) -> Dict[str, Any]:
    logging.info(f"Processing TXT: {txt_path.name}")
    with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
        raw_text = f.read()
    text = normalize_text(raw_text)
    words = text.split()
    parts = split_with_overlap(words, cfg.get("PARA_TARGET_TOKENS", 300), cfg.get("PARA_OVERLAP_TOKENS", 50))
    dedup_intensity = cfg.get("DEDUPLICATION_INTENSITY", 1)
    seen_hashes = {}
    detail_records = []
    for idx, part in enumerate(parts):
        if not should_deduplicate(part, seen_hashes, dedup_intensity):
            detail_records.append({
                "id": f"{txt_path.name}::chunk{idx+1:03d}::{sha8(part)}",
                "text": part, "text_raw": part, "element_type": "paragraph",
                "metadata": {"doc_id": txt_path.name, "chapter_title": txt_path.stem,
                             "chapter_id": f"chunk{idx+1:03d}", "breadcrumb": txt_path.stem,
                             "hierarchy_level": 1, "page_start": None, "page_end": None},
                "raw_markdown": None
            })
    summary = summarize_for_router(text, cfg.get("MAX_ROUTER_SUMMARY_CHARS", 3000))
    return {
        "router": [{"route_id": f"{txt_path.name}::doc", "title": txt_path.stem,
                     "breadcrumb": txt_path.stem, "level": 1, "scope_pages": [1, 1],
                     "summary": summary, "tags": []}],
        "detail": detail_records
    }


def build_for_docx(docx_path: Path, cfg: Dict[str, Any]) -> Dict[str, Any]:
    logging.info(f"Processing DOCX: {docx_path.name}")
    from docx import Document
    doc = Document(str(docx_path))
    image_texts = process_docx_images(doc, cfg)
    ocr_text = "\n\n[OCR from images]\n" + "\n\n".join([img["text"] for img in image_texts]) if image_texts else ""
    sections = []
    current_section = {"title": docx_path.stem, "level": 1, "text": []}
    for para in doc.paragraphs:
        text = normalize_text(para.text)
        if not text:
            continue
        if para.style.name.startswith('Heading'):
            if current_section["text"]:
                sections.append(current_section)
            level = int(para.style.name.replace('Heading', '').strip() or '1')
            current_section = {"title": text, "level": level, "text": []}
        else:
            current_section["text"].append(text)
    if current_section["text"]:
        sections.append(current_section)
    if ocr_text and sections:
        sections[-1]["text"].append(ocr_text)
    if not sections:
        all_text = " ".join([normalize_text(p.text) for p in doc.paragraphs])
        sections = [{"title": docx_path.stem, "level": 1, "text": [all_text]}]

    router_records, detail_records = [], []
    dedup_intensity = cfg.get("DEDUPLICATION_INTENSITY", 1)
    seen_hashes = {}
    target_tokens = cfg.get("PARA_TARGET_TOKENS", 300)
    overlap_tokens = cfg.get("PARA_OVERLAP_TOKENS", 50)
    for idx, section in enumerate(sections):
        section_text = " ".join(section["text"])
        words = section_text.split()
        parts = split_with_overlap(words, target_tokens, overlap_tokens)
        ch_id = f"sec{idx+1:02d}"
        breadcrumb = section["title"]
        router_records.append({
            "route_id": f"{docx_path.name}::{ch_id}", "title": section["title"],
            "breadcrumb": breadcrumb, "level": section["level"],
            "scope_pages": [idx+1, idx+1],
            "summary": summarize_for_router(section_text, cfg.get("MAX_ROUTER_SUMMARY_CHARS", 3000)),
            "tags": []
        })
        for part_idx, part in enumerate(parts):
            if not should_deduplicate(part, seen_hashes, dedup_intensity):
                detail_records.append({
                    "id": f"{docx_path.name}::{ch_id}::part{part_idx+1:03d}::{sha8(part)}",
                    "text": f"[{breadcrumb}]\n{part}", "text_raw": part,
                    "element_type": "paragraph",
                    "metadata": {"doc_id": docx_path.name, "chapter_title": section["title"],
                                 "chapter_id": ch_id, "breadcrumb": breadcrumb,
                                 "hierarchy_level": section["level"],
                                 "page_start": idx+1, "page_end": idx+1},
                    "raw_markdown": None
                })
    return {"router": router_records, "detail": detail_records}


def build_for_pptx(pptx_path: Path, cfg: Dict[str, Any]) -> Dict[str, Any]:
    logging.info(f"Processing PPTX: {pptx_path.name}")
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    image_texts = process_pptx_images(prs, cfg)
    ocr_by_slide = {}
    for img_data in image_texts:
        ocr_by_slide.setdefault(img_data["slide"], []).append(img_data["text"])
    router_records, detail_records = [], []
    dedup_intensity = cfg.get("DEDUPLICATION_INTENSITY", 1)
    seen_hashes = {}
    target_tokens = cfg.get("PARA_TARGET_TOKENS", 300)
    overlap_tokens = cfg.get("PARA_OVERLAP_TOKENS", 50)
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if getattr(slide, "shapes", None) and slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        bullets = []
        for shp in slide.shapes:
            if hasattr(shp, "text") and shp.text:
                bullets.append(shp.text.strip())
        text = " ".join(bullets).strip()
        if i in ocr_by_slide:
            text += "\n\n[OCR from images]\n" + "\n\n".join(ocr_by_slide[i])
        router_records.append({
            "route_id": f"{pptx_path.name}::slide{i:03d}",
            "title": title or f"Slide {i}", "scope_pages": [i, i],
            "summary": text[:900], "tags": ["slides"]
        })
        words = text.split()
        parts = split_with_overlap(words, target_tokens, overlap_tokens) if words else []
        for part_idx, part in enumerate(parts):
            if not should_deduplicate(part, seen_hashes, dedup_intensity):
                detail_records.append({
                    "id": f"{pptx_path.name}::slide{i:03d}::part{part_idx+1:03d}::{sha8(part)}",
                    "text": part, "text_raw": part, "element_type": "paragraph",
                    "metadata": {"doc_id": pptx_path.name, "chapter_id": f"slide{i:03d}",
                                 "page_start": i, "page_end": i},
                    "raw_markdown": None
                })
    return {"router": router_records, "detail": detail_records}
